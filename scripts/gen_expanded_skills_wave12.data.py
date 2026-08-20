SKILLS = [
    {
        "name": "research-paper-ideation",
        "title": "Research Paper Ideation with AI",
        "description": "Use LLMs, citation networks, and structured brainstorming to generate and refine research questions, hypotheses, and project outlines.",
        "devin_body": r'''
## When to use

You are starting a new research project, exploring a new domain, or need to turn a broad interest into a focused, novel research question.

## Key concepts

- **Research question trees**: decompose a broad topic into nested, testable questions.
- **Literature gap analysis**: identify what has not been done by mapping existing work.
- **Hypothesis generation**: form falsifiable claims that connect methods, data, and outcomes.
- **Concept mapping**: visualize relationships between variables, mechanisms, and prior findings.
- **AI ideation agents**: use retrieval-augmented LLMs to propose, refine, and evaluate ideas against real papers.

## Code pattern

```python
import requests
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans


def fetch_related_papers(paper_id, fields="title,abstract,year"):
    url = f"https://api.semanticscholar.org/graph/v1/paper/{paper_id}/references"
    params = {"fields": fields, "limit": 100}
    return requests.get(url, params=params, timeout=30).json()["data"]


def cluster_papers(papers, n=5):
    texts = [
        p["paper"]["title"] + " " + (p["paper"]["abstract"] or "")
        for p in papers
    ]
    X = TfidfVectorizer(stop_words="english", max_features=500).fit_transform(texts)
    labels = KMeans(n_clusters=n, random_state=42, n_init="auto").fit_predict(X)
    return pd.DataFrame(
        {"title": [p["paper"]["title"] for p in papers], "cluster": labels}
    )
```

## Tuning notes

- Always ground AI suggestions in real literature; verify citations exist and are relevant.
- Distinguish novelty from feasibility; a novel but infeasible idea is not a good project.
- Use multiple ideation rounds and diverse prompts to avoid anchoring on the first idea.
- Involve collaborators early to challenge assumptions and sharpen the question.

## Verification

1. Generate 5 candidate research questions for a target domain.
2. Map at least 30 related papers and cluster them to find gaps.
3. Pick one question and outline hypotheses, methods, and expected outcomes.
''',
        "references": [
            "https://doi.org/10.1016/j.ijresmar.2023.10.002",
            "https://arxiv.org/abs/2503.00946v3",
            "https://arxiv.org/pdf/2409.04109",
            "https://journals.plos.org/ploscompbiol/article?id=10.1371/journal.pcbi.1013259",
            "https://openreview.net/forum?id=bIAFQ8asqi",
        ],
    },
    {
        "name": "grant-proposal-writing",
        "title": "Grant Proposal Writing with AI",
        "description": "Structure Specific Aims, research strategy, budget, and broader impact sections for NIH/NSF/ERC-style proposals with AI drafting support.",
        "devin_body": r'''
## When to use

You are preparing a fellowship, R01, NSF CAREER, Horizon Europe, or other competitive research proposal and need a clear, compelling narrative.

## Key concepts

- **Specific Aims**: a concise statement of goals, hypotheses, and expected outcomes.
- **Significance / Innovation / Approach**: the core review criteria for many funders.
- **Broader impacts**: training, dissemination, societal benefit, and reproducibility.
- **Budget justification**: link personnel, equipment, and travel directly to aims.
- **Funder compliance**: follow page limits, formatting, and required sections exactly.

## Code pattern

```python
import yaml
from datetime import datetime, timedelta


def build_proposal_outline(title, aims, duration_months=36):
    outline = {
        "title": title,
        "specific_aims": aims,
        "research_strategy": ["significance", "innovation", "approach"],
        "timeline": [
            (datetime.now() + timedelta(days=30 * i)).strftime("%Y-%m")
            for i in range(duration_months)
        ],
        "broader_impacts": [],
        "budget_justification": [],
    }
    with open("proposal_outline.yaml", "w") as f:
        yaml.safe_dump(outline, f, sort_keys=False)
    return outline
```

## Tuning notes

- Write for both expert and non-expert reviewers; clarity beats complexity.
- Align every aim with a budget line and a measurable milestone.
- Use figures and timelines to make the approach concrete.
- Avoid over-promising; funders value realistic, well-scoped work plans.

## Verification

1. Draft a one-page Specific Aims document and check it against funder guidelines.
2. Build a full proposal outline and cross-reference every section to the instructions.
3. Share with a mentor or prior awardee for feedback on significance and feasibility.
''',
        "references": [
            "https://www.nigms.nih.gov/Research/application/Pages/default",
            "https://www.nigms.nih.gov/training/Pages/Grant-Writing-Webinar-Series-for-Institutions-Building-Research--and-Research-Training-Capacity",
            "https://www.nigms.nih.gov/Research/application/Pages/Submitting-an-Application",
            "https://www.nimh.nih.gov/funding/grant-writing-and-application-process/grant-writing-assistance",
            "https://blogs.nature.com/blog/beginnings-how-to-write-your-first-grant-proposal/",
        ],
    },
    {
        "name": "scientific-writing",
        "title": "Scientific Writing with AI",
        "description": "Improve clarity, structure, and style for manuscripts, theses, and reports using AI drafting and editing tools.",
        "devin_body": r'''
## When to use

You are drafting a manuscript, revising for a journal, or trying to make complex research accessible to a broader audience.

## Key concepts

- **IMRAD structure**: Introduction, Methods, Results, And Discussion.
- **C-C-C scheme**: Context-Content-Conclusion at the paragraph level.
- **Active voice and parallel structure**: improve readability and momentum.
- **Readability metrics**: Flesch-Kincaid grade, sentence length, word complexity.
- **Central contribution**: every section should reinforce the paper's main message.

## Code pattern

```python
import textstat
import re


def analyze_readability(text):
    sentences = [s for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]
    word_counts = [len(s.split()) for s in sentences]
    return {
        "flesch_kincaid_grade": textstat.flesch_kincaid_grade(text),
        "avg_sentence_length": sum(word_counts) / len(word_counts),
        "word_count": len(text.split()),
        "sentence_count": len(sentences),
    }


# Example: load a draft and flag overly long sentences
with open("draft.txt") as f:
    report = analyze_readability(f.read())
print(report)
```

## Tuning notes

- Put the central contribution in the title, abstract, and first paragraph.
- Write for flesh-and-blood readers who do not already know your work.
- Use AI for revision, not for fabricating citations or results.
- Verify that every claim in the introduction is supported in the results.

## Verification

1. Analyze a draft for readability and sentence length.
2. Restructure one section using the C-C-C scheme.
3. Compare the before and after versions with a co-author or reader.
''',
        "references": [
            "https://doi.org/10.1371/journal.pcbi.1003453",
            "https://journals.plos.org/ploscompbiol/article?id=10.1371/journal.pcbi.1005619",
            "https://www.gatsby.ucl.ac.uk/~pel/misc/gopen_swan.pdf",
            "https://www.coursera.org/learn/sciwrite",
        ],
    },
    {
        "name": "ai-peer-review",
        "title": "AI-Assisted Peer Review",
        "description": "Use AI tools and structured checklists to write constructive, ethical peer reviews for manuscripts and proposals.",
        "devin_body": r'''
## When to use

You are reviewing a manuscript, preprint, or conference submission and want to produce a fair, structured, and actionable review.

## Key concepts

- **COPE guidelines**: ethical standards for reviewers, editors, and authors.
- **Novelty, significance, rigor**: core dimensions of scientific evaluation.
- **Confidentiality and conflict of interest**: protect unpublished work and declare biases.
- **Constructive critique**: separate major concerns from minor suggestions.
- **Responsible AI use**: disclose any AI assistance and verify generated claims.

## Code pattern

```python
from pathlib import Path


def review_checklist(manuscript_text):
    checklist = {
        "novelty_and_significance": False,
        "methods_and_rigor": False,
        "data_availability": False,
        "ethical_approval": False,
        "conflict_of_interest": False,
        "ai_assistance_disclosed": False,
        "constructive_tone": False,
    }
    # Use keyword checks as a starting point; human judgment is required.
    if "data availability" in manuscript_text.lower():
        checklist["data_availability"] = True
    if "ethical" in manuscript_text.lower() or "irb" in manuscript_text.lower():
        checklist["ethical_approval"] = True
    return checklist


with open("manuscript.txt") as f:
    print(review_checklist(f.read()))
```

## Tuning notes

- Never upload a confidential manuscript into a public AI tool without permission.
- Use AI to organize your own notes, not to replace domain judgment.
- Be specific: cite line numbers, figures, or equations when raising issues.
- Distinguish required revisions from optional suggestions.

## Verification

1. Write a review of a sample paper using the checklist.
2. Compare your review against the journal's reviewer guidelines.
3. Have a colleague read it and confirm the tone is constructive and fair.
''',
        "references": [
            "https://www.nature.com/nature/for-referees",
            "https://www.nature.com/nature/for-referees/how-to-write-a-report",
            "https://www.nature.com/nm/editorial-policies/ai",
            "https://journals.plos.org/ploscompbiol/article?id=10.1371/journal.pcbi.0020110",
            "https://www.nature.com/articles/s41565-026-02177-2",
        ],
    },
    {
        "name": "citation-management",
        "title": "Citation Management",
        "description": "Organize references, manage PDFs, format bibliographies, and share libraries with Zotero, Mendeley, or BibTeX.",
        "devin_body": r'''
## When to use

You are building a literature library, collaborating on a manuscript, or switching between citation styles for different venues.

## Key concepts

- **Reference manager**: Zotero, Mendeley, EndNote, JabRef, or BibTeX-based tools.
- **Citation Style Language (CSL)**: format bibliographies in thousands of styles.
- **Metadata cleanup**: verify DOIs, author names, journal titles, and page numbers.
- **Group libraries and shared collections**: collaborate with co-authors.
- **Import/export formats**: RIS, BibTeX, CSL-JSON, and Zotero connectors.

## Code pattern

```python
import bibtexparser
from collections import Counter


def load_and_dedup(bib_path):
    with open(bib_path) as f:
        db = bibtexparser.load(f)
    seen = set()
    unique = []
    for entry in db.entries:
        key = (entry.get("doi") or "").lower() or (
            entry.get("title", "") + entry.get("year", "")
        )
        if key not in seen:
            seen.add(key)
            unique.append(entry)
    db.entries = unique
    return db


def style_counts(db):
    return Counter(entry.get("ENTRYTYPE", "unknown") for entry in db.entries)
```

## Tuning notes

- Always verify imported metadata; PDF metadata is often noisy or incomplete.
- Use DOIs as stable identifiers and link them to Crossref for updates.
- Back up your library and sync across devices.
- Keep one master library and create project-specific collections.

## Verification

1. Import 20 references and generate bibliographies in APA and Vancouver styles.
2. Identify and merge duplicate entries.
3. Verify that all citations resolve to real DOIs or URLs.
''',
        "references": [
            "https://www.zotero.org/support/quick_start_guide/",
            "https://www.zotero.org/support/styles",
            "https://service.elsevier.com/app/answers/detail/a_id/29356/supporthub/mendeley/role/referencemanagement/",
            "https://doi.org/10.1371/journal.pcbi.1006036",
            "https://doi.org/10.1038/npre.2009.3867.1",
        ],
    },
    {
        "name": "research-presentation-design",
        "title": "Research Presentation Design",
        "description": "Build clear, compelling slides and posters for seminars, conferences, and outreach using narrative structure and visual hierarchy.",
        "devin_body": r'''
## When to use

You are preparing a conference talk, seminar, thesis defense, poster session, or public outreach presentation.

## Key concepts

- **One idea per slide**: keep each slide focused on a single message.
- **Assertion-Evidence structure**: state a claim and support it with a visual.
- **Visual hierarchy**: use size, color, and position to guide attention.
- **Data-ink ratio**: remove unnecessary chart elements and decoration.
- **Accessibility**: use readable fonts, high contrast, and alt text.

## Code pattern

```python
from pptx import Presentation
from pptx.util import Inches, Pt


prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
)
tf = title_box.text_frame
tf.text = "One clear message per slide"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True

# Add a placeholder for a figure
left = Inches(1)
top = Inches(1.8)
slide.shapes.add_picture("plot.png", left, top, height=Inches(4.5))
prs.save("presentation.pptx")
```

## Tuning notes

- Design for the back of the room: large fonts, simple figures, minimal text.
- Practice timing; a 15-minute talk needs a tight narrative arc.
- Use consistent colors, fonts, and alignment across all slides.
- Include a take-home slide with the main conclusion and contact info.

## Verification

1. Build a 10-slide deck and test it on a projector or large screen.
2. Critique a peer's slides for visual hierarchy and data-ink ratio.
3. Deliver the talk to a practice audience and gather feedback.
''',
        "references": [
            "https://doi.org/10.1371/journal.pcbi.0030077",
            "https://journals.plos.org/ploscompbiol/article?id=10.1371/journal.pcbi.1009554",
            "https://journals.plos.org/ploscompbiol/article?id=10.1371/journal.pcbi.1005373",
            "https://journals.plos.org/ploscompbiol/article?id=10.1371/journal.pcbi.1007163",
        ],
    },
    {
        "name": "market-research-ai",
        "title": "AI for Market Research",
        "description": "Design surveys, segment customers, analyze open-ended responses, and forecast market trends with AI-driven tools.",
        "devin_body": r'''
## When to use

You need to assess product-market fit, customer preferences, pricing sensitivity, or competitive positioning for a product or research spin-out.

## Key concepts

- **Survey design**: clear questions, response scales, sampling, and bias control.
- **Conjoint and MaxDiff**: measure feature and price preferences.
- **Synthetic panels and LLM responses**: fast, low-cost but require validation.
- **Text analysis of open-ends**: topic modeling, sentiment, and theme extraction.
- **Trend forecasting**: time-series models and leading-indicator tracking.

## Code pattern

```python
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import NMF


def analyze_open_ends(responses, n_topics=5):
    vectorizer = TfidfVectorizer(stop_words="english", max_features=500)
    X = vectorizer.fit_transform(responses)
    nmf = NMF(n_components=n_topics, random_state=42, max_iter=500)
    W = nmf.fit_transform(X)
    terms = vectorizer.get_feature_names_out()
    topics = [
        [terms[i] for i in topic.argsort()[-5:]]
        for topic in nmf.components_
    ]
    scores = pd.DataFrame(
        W, columns=[f"topic_{i}" for i in range(n_topics)]
    )
    return scores, topics
```

## Tuning notes

- Pre-test survey questions to avoid ambiguity and leading wording.
- Validate synthetic or LLM-generated responses against a small real panel.
- Protect respondent privacy and comply with data-use agreements.
- Combine quantitative scores with qualitative quotes for richer insight.

## Verification

1. Design a short survey for a product concept.
2. Collect or simulate 50+ responses and analyze the open-ended answers.
3. Compare AI-derived themes with a manual coding of a subset.
''',
        "references": [
            "https://www.qualtrics.com/articles/strategy-research/agentic-ai-market-research/",
            "https://www.hbs.edu/ris/Publication%20Files/23-062_1f58623a-ee21-44b9-a262-276047bc5543.pdf",
            "https://www.surveymonkey.com/use-cases/market-research/",
            "https://esocorpwebsitestg.blob.core.windows.net/strapi-uploads/uploads/cltn6755401khqe3v0od2y6ut_esomar_20_questions_to_help_buyers_of_ai_based_services_0277e1b5eb.pdf",
        ],
    },
    {
        "name": "competitive-analysis",
        "title": "Competitive Analysis with AI",
        "description": "Map industry structure, benchmark competitors, and identify strategic positioning using Porter's Five Forces, SWOT, and data.",
        "devin_body": r'''
## When to use

You are entering a new market, launching a product, writing an industry background section, or planning a strategic pivot.

## Key concepts

- **Porter's Five Forces**: rivalry, new entrants, substitutes, buyer power, supplier power.
- **SWOT**: Strengths, Weaknesses, Opportunities, Threats.
- **SCP framework**: Structure-Conduct-Performance and industry attractiveness.
- **Competitor profiling**: products, pricing, positioning, and capabilities.
- **Strategic positioning**: where to play and how to win.

## Code pattern

```python
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt


def five_forces_radar(scores):
    categories = list(scores.keys())
    values = list(scores.values())
    angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
    values += values[:1]
    angles += angles[:1]

    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    ax.plot(angles, values, "o-", linewidth=2)
    ax.fill(angles, values, alpha=0.25)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories)
    ax.set_ylim(0, 5)
    plt.title("Five Forces Scorecard")
    plt.savefig("five_forces_radar.png", dpi=300)


five_forces_radar({
    "rivalry": 4,
    "new_entrants": 2,
    "substitutes": 3,
    "buyer_power": 3,
    "supplier_power": 2,
})
```

## Tuning notes

- Ground scores in public data and primary sources, not intuition alone.
- Update the analysis regularly; competitive landscapes shift quickly.
- Distinguish facts from inference and label assumptions clearly.
- Pair Five Forces with a SWOT to capture internal and external factors.

## Verification

1. Build a Five Forces scorecard and a SWOT grid for a target market.
2. Compare your analysis with a published industry report.
3. Present it to stakeholders and test whether it informs decisions.
''',
        "references": [
            "https://hbr.org/1979/03/how-competitive-forces-shape-strategy",
            "https://hbr.org/2008/01/the-five-competitive-forces-that-shape-strategy",
            "https://hbr.org/2021/02/are-you-doing-the-swot-analysis-backwards",
            "https://www.coursera.org/articles/competitor-analysis",
            "https://www.sba.gov/business-guide/plan-your-business/market-research-competitive-analysis",
        ],
    },
    {
        "name": "product-requirements-ai",
        "title": "AI for Product Requirements",
        "description": "Draft, validate, and track product requirements documents (PRDs) with user stories, assumptions, and success metrics.",
        "devin_body": r'''
## When to use

You are scoping a new feature or product, aligning engineering and design, or need a single source of truth for what to build and why.

## Key concepts

- **PRD (Product Requirements Document)**: defines purpose, features, behavior, and success criteria.
- **User stories and acceptance criteria**: who, what, and why for each capability.
- **Success metrics**: measurable outcomes tied to user and business goals.
- **Assumptions and out-of-scope**: manage risk and prevent scope creep.
- **Prioritization**: rank features by value, effort, and strategic fit.

## Code pattern

```python
import yaml
from datetime import date


def draft_prd(problem, solution, audience, metrics):
    prd = {
        "problem": problem,
        "solution_vision": solution,
        "target_audience": audience,
        "success_metrics": metrics,
        "assumptions": [],
        "out_of_scope": [],
        "user_stories": [],
        "last_updated": date.today().isoformat(),
    }
    with open("prd.yaml", "w") as f:
        yaml.safe_dump(prd, f, sort_keys=False)
    return prd


def add_user_story(prd, role, want, so_that, acceptance=[]):
    prd["user_stories"].append({
        "role": role,
        "want": want,
        "so_that": so_that,
        "acceptance_criteria": acceptance,
    })
    return prd
```

## Tuning notes

- Start with the problem and the user, not the implementation.
- Keep PRDs concise and living; update them as discovery progresses.
- Align engineering, design, and stakeholders before writing detailed specs.
- Tie each requirement to a success metric or user outcome.

## Verification

1. Write a one-page PRD for a feature.
2. Run it by an engineer, a designer, and a potential user.
3. Trace each user story to at least one design mock-up and one test.
''',
        "references": [
            "https://www.atlassian.com/agile/product-management/requirements",
            "https://www.atlassian.com/software/confluence/templates/product-requirements",
            "https://www.atlassian.com/software/confluence/templates/requirements",
            "https://confluence.atlassian.com/doc/blog/2015/08/how-to-document-product-requirements-in-confluence",
            "https://www.svpg.com/wp-content/uploads/2024/07/How-To-Write-a-Good-PRD.pdf",
        ],
    },
    {
        "name": "user-interviews-synthesis",
        "title": "User Interview Synthesis",
        "description": "Turn interview transcripts into themes, insights, and personas using thematic analysis, affinity mapping, and AI coding.",
        "devin_body": r'''
## When to use

You have completed qualitative user or stakeholder interviews and need to extract actionable themes, insights, and design implications.

## Key concepts

- **Thematic analysis**: identify, analyze, and report patterns across data.
- **Open and axial coding**: label quotes and group codes into themes.
- **Affinity diagramming**: cluster observations collaboratively on sticky notes or digital boards.
- **Personas and empathy maps**: synthesize user needs and contexts.
- **Saturation**: stop sampling when new interviews add no new themes.

## Code pattern

```python
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import NMF


def code_transcripts(csv_path, n_themes=5):
    df = pd.read_csv(csv_path)
    quotes = df["quote"].dropna().astype(str)
    vectorizer = TfidfVectorizer(
        stop_words="english", ngram_range=(1, 2), max_features=500
    )
    X = vectorizer.fit_transform(quotes)
    nmf = NMF(n_components=n_themes, random_state=42, max_iter=500)
    W = nmf.fit_transform(X)
    df["dominant_theme"] = W.argmax(axis=1)
    terms = vectorizer.get_feature_names_out()
    themes = [
        [terms[i] for i in comp.argsort()[-5:]]
        for comp in nmf.components_
    ]
    return df, themes


df, themes = code_transcripts("interview_quotes.csv")
print(themes)
```

## Tuning notes

- Start synthesis with clear research questions, not the algorithm.
- Avoid confirmation bias; seek disconfirming evidence and edge cases.
- Involve the team; interpretation benefits from multiple perspectives.
- Triangulate interview findings with surveys, analytics, or prototypes.

## Verification

1. Synthesize 3-5 transcripts and produce a theme report.
2. Build an affinity diagram and compare it to the algorithmic output.
3. Compare theme assignments with an independent rater and compute agreement.
''',
        "references": [
            "https://www.nngroup.com/articles/affinity-diagram/",
            "https://www.userinterviews.com/blog/affinity-mapping-ux-research-data-synthesis",
            "https://dovetail.com/research/research-synthesis/",
            "https://handbook.gitlab.com/handbook/upstream-studios/experience-research/analyzing-research-data/",
            "https://www2.uwe.ac.uk/services/Marketing/students/Newstudents/HAS/Using%20thematic%20analysis%20in%20psychology.pdf",
        ],
    },
    {
        "name": "research-data-storytelling",
        "title": "Research Data Storytelling",
        "description": "Turn complex scientific results into narrative visualizations and stories that resonate with specialists and the public.",
        "devin_body": r'''
## When to use

You are preparing figures, a press release, a public talk, a grant impact section, or any communication where the story behind the data matters.

## Key concepts

- **Narrative arc**: setup, tension, resolution, and call to action.
- **Audience-centric design**: match the message to the reader's background.
- **Data-ink ratio and chart junk**: remove non-essential marks.
- **Annotations and callouts**: guide attention to the key data point.
- **Ethical representation**: show uncertainty, avoid misleading axes, and credit sources.

## Code pattern

```python
import matplotlib.pyplot as plt


def story_chart(years, values, highlight_year, message):
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.plot(years, values, marker="o", color="steelblue")

    idx = years.index(highlight_year)
    ax.annotate(
        message,
        xy=(highlight_year, values[idx]),
        xytext=(highlight_year + 1, values[idx] + 0.5),
        arrowprops=dict(facecolor="black", shrink=0.05, width=1, headwidth=6),
    )
    ax.set_title(message, fontsize=14, weight="bold")
    ax.set_ylabel("Measured value")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig("story_figure.png", dpi=300)


story_chart(
    years=[2020, 2021, 2022, 2023, 2024],
    values=[12, 15, 18, 22, 35],
    highlight_year=2024,
    message="Treatment effect doubled after 2023",
)
```

## Tuning notes

- Identify one clear message for each visual and build around it.
- Use plain language titles; avoid jargon in public-facing graphics.
- Include error bars, confidence intervals, or sample sizes where relevant.
- Test visuals with a non-expert to ensure the story is clear.

## Verification

1. Design a figure that tells a specific story from your data.
2. Gather feedback from both a specialist and a non-specialist.
3. Check accessibility: color-blind safe palette, alt text, and readable labels.
''',
        "references": [
            "https://doi.org/10.1109/tvcg.2010.179",
            "https://doi.org/10.1371/journal.pcbi.1003833",
            "https://help.tableau.com/current/pro/desktop/en-gb/story_best_practices.htm",
            "https://www.storytellingwithdata.com/books",
            "https://books.google.com/books/about/Storytelling_with_Data.html?id=rRSRCgAAQBAJ",
        ],
    },
    {
        "name": "collaboration-and-team-science",
        "title": "Collaboration and Team Science",
        "description": "Build, lead, and sustain productive interdisciplinary research teams with clear roles, communication, and shared tools.",
        "devin_body": r'''
## When to use

You are assembling a research team, coordinating a multi-site or interdisciplinary project, or establishing authorship and data-sharing agreements.

## Key concepts

- **Team science**: collaborative, often interdisciplinary research to tackle complex problems.
- **Role clarity**: define who does what, including data, methods, writing, and management.
- **Psychological safety**: create an environment where team members can raise concerns.
- **Authorship and contribution agreements**: decide early and revisit regularly.
- **Communication cadence and shared infrastructure**: version control, shared drives, and meeting rituals.

## Code pattern

```python
import csv
from datetime import datetime


def create_team_charter(roles, contributions):
    with open("team_charter.csv", "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["name", "role", "contribution", "start_date"])
        for member in roles:
            writer.writerow([
                member["name"],
                member["role"],
                contributions.get(member["name"], ""),
                datetime.now().isoformat(),
            ])


def authorship_matrix(members, tasks):
    matrix = {m: {t: 0 for t in tasks} for m in members}
    # Populate manually or from contribution logs
    return matrix
```

## Tuning notes

- Agree on goals, roles, authorship, and data-sharing rules at kickoff.
- Schedule regular syncs and keep decision logs to reduce misalignment.
- Use shared repositories and documents; avoid siloed files and email chains.
- Address conflict early and revisit team norms at major milestones.

## Verification

1. Create a team charter with roles, contributions, and authorship principles.
2. Publish a shared data and authorship plan that all members approve.
3. Review the charter and collaboration health at each project milestone.
''',
        "references": [
            "https://www.cancer.gov/about-nci/organization/crs/research-initiatives/team-science-field-guide/collaboration-team-science-guide.pdf?rid=267&tid=1",
            "https://nap.nationalacademies.org/resource/29043/interactive/",
            "https://www.ncbi.nlm.nih.gov/books/NBK617881/",
            "https://pmc.ncbi.nlm.nih.gov/articles/PMC8599160/",
            "https://www.nigms.nih.gov/grants/Pages/Considerations-for-Multiple-Principal-Investigator-Applications",
        ],
    },
]
